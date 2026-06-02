# -*- coding: utf-8 -*-
"""YouTube 채널 관심도 분석기 — 발표 PPT 생성 스크립트"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

IMG = os.path.join(os.path.dirname(__file__), "static", "images")

# ---- 테마 색상 (앱 다크 테마와 동일 계열) ----
BG     = RGBColor(0x0A, 0x0A, 0x0B)   # 거의 검정
PANEL  = RGBColor(0x16, 0x16, 0x19)   # 카드 패널
RED    = RGBColor(0xFF, 0x3B, 0x3B)   # 강조 빨강
GREEN  = RGBColor(0x2E, 0xCC, 0x71)   # 연결됨 초록
ORANGE = RGBColor(0xF4, 0xA8, 0x44)   # 보류/태그
WHITE  = RGBColor(0xF2, 0xF2, 0xF2)
GREY   = RGBColor(0x9A, 0x9A, 0xA2)
DIM    = RGBColor(0x6A, 0x6A, 0x72)

FONT = "Malgun Gothic"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # 좌측 빨강 액센트 바
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.13), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background(); bar.shadow.inherit = False
    return s


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = line_spacing
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb


def R(txt, size, color=WHITE, bold=False, font=FONT):
    return (txt, size, color, bold, font)


def title(s, kicker, head, headcolor=WHITE):
    text(s, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.4),
         [[R(kicker, 13, RED, True, MONO)]])
    text(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.7),
         [[R(head, 30, headcolor, True)]])
    box(s, Inches(0.57), Inches(1.32), Inches(2.1), Pt(2.5), fill=RED)


def add_image_fit(s, path, l, t, w, h, border=True):
    """비율 유지하며 (l,t,w,h) 박스 안에 가운데 배치."""
    im = Image.open(path); iw, ih = im.size
    bw, bh = float(w), float(h)
    scale = min(bw / iw, bh / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    nl = int(l) + int((bw - nw) / 2)
    nt = int(t) + int((bh - nh) / 2)
    pic = s.shapes.add_picture(path, nl, nt, Emu(nw), Emu(nh))
    if border:
        ln = pic.line; ln.color.rgb = RGBColor(0x33, 0x33, 0x3A); ln.width = Pt(1)
    return pic


def imgpath(name):
    return os.path.join(IMG, name)


# ============================================================
# 슬라이드 1 — 타이틀
# ============================================================
s = slide()
box(s, Inches(0), Inches(0), SW, SH, fill=BG)
box(s, Inches(0), Inches(0), Inches(0.13), SH, fill=RED)
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     [[R("// YOUTUBE CHANNEL ANALYZER", 18, RED, True, MONO)]])
text(s, Inches(0.85), Inches(2.35), Inches(11.7), Inches(1.8),
     [[R("구독 채널 ", 54, WHITE, True), R("관심도 분석기", 54, RED, True)]])
text(s, Inches(0.9), Inches(3.95), Inches(11.0), Inches(1.2),
     [[R("Google OAuth로 내 계정을 연결해 구독 채널을 ", 18, GREY),
       R("관심도 점수(0~100)", 18, WHITE, True), R("로 환산하고,", 18, GREY)],
      [R("Solar LLM이 시청 취향을 해석해 ", 18, GREY),
       R("페르소나·구독 정리·새 채널", 18, WHITE, True),
       R("까지 추천하는 웹 앱", 18, GREY)]], line_spacing=1.15)
# 기술 칩
chips = ["Flask", "Vanilla JS", "Google OAuth (GIS)", "YouTube Data API v3", "Upstage Solar LLM"]
cx = Inches(0.9)
for c in chips:
    w = Inches(0.35 + 0.115 * len(c))
    b = box(s, cx, Inches(5.45), w, Inches(0.46), fill=PANEL,
            line=RGBColor(0x3A, 0x3A, 0x42), radius=True)
    text(s, cx, Inches(5.45), w, Inches(0.46), [[R(c, 12, WHITE, True, MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx = Emu(int(cx) + int(w) + Inches(0.18))
text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
     [[R("알고리즘 실습 — 팀 프로젝트 발표", 14, DIM, False, MONO)]])

# ============================================================
# 슬라이드 2 — 프로젝트 개요 / 문제의식
# ============================================================
s = slide()
title(s, "01 · OVERVIEW", "이런 문제에서 출발했습니다")
# 좌: 문제
box(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(5.3), fill=PANEL, radius=True)
text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.5),
     [[R("문제의식", 18, RED, True)]])
text(s, Inches(0.85), Inches(2.45), Inches(5.45), Inches(4.2),
     [[R("• 구독 채널이 수십~수백 개로 쌓이지만", 15, WHITE)],
      [R("  ", 15, WHITE), R("정작 어떤 채널을 진짜로 보는지", 15, GREY)],
      [R("  알기 어렵다.", 15, GREY)],
      [R("", 8, GREY)],
      [R("• \"10년 구독했지만 안 보는\" 휴면 채널이", 15, WHITE)],
      [R("  목록을 가득 채우고 있다.", 15, GREY)],
      [R("", 8, GREY)],
      [R("• YouTube는 \"내 댓글·시청 기록 전체\"를", 15, WHITE)],
      [R("  조회하는 API를 제공하지 않는다", 15, GREY),
       R(" (2016 폐지).", 13, DIM)]], line_spacing=1.1)
# 우: 해결
box(s, Inches(6.78), Inches(1.6), Inches(6.0), Inches(5.3), fill=PANEL, radius=True)
text(s, Inches(7.08), Inches(1.85), Inches(5.5), Inches(0.5),
     [[R("우리의 해결", 18, GREEN, True)]])
text(s, Inches(7.08), Inches(2.45), Inches(5.45), Inches(4.2),
     [[R("✓ 댓글·좋아요·시청 횟수·구독 기간을", 15, WHITE)],
      [R("   가중 합산 → ", 15, GREY), R("관심도 점수(0~100)", 15, WHITE, True)],
      [R("", 8, GREY)],
      [R("✓ ", 15, WHITE), R("활동 0이면 즉시 0점", 15, RED, True),
       R(" — 휴면 강등", 15, GREY)],
      [R("", 8, GREY)],
      [R("✓ 4그룹 자동 분류 🔥 😊 👋 🌱", 15, WHITE)],
      [R("", 8, GREY)],
      [R("✓ Google Takeout으로 ", 15, WHITE),
       R("실제 시청 기록", 15, WHITE, True), R(" 보강", 15, GREY)],
      [R("", 8, GREY)],
      [R("✓ Solar LLM이 취향 해석·채널 추천", 15, WHITE)]], line_spacing=1.1)

# ============================================================
# 슬라이드 3 — 핵심 기능
# ============================================================
s = slide()
title(s, "02 · FEATURES", "핵심 기능 한눈에")
feats = [
    ("📊", "관심도 점수화", "댓글·좋아요·시청·구독 기간을\n가중 합산해 0~100 점으로 환산"),
    ("🗂️", "4그룹 자동 분류", "🔥 최애 / 😊 자주 / 👋 가끔 /\n🌱 구독만 — 이진 탐색으로 매핑"),
    ("📥", "Takeout 하이브리드", "comments.csv·watch-history.json\n으로 정확한 댓글·시청 기록 반영"),
    ("🧬", "시청 취향 페르소나", "채널 분포를 Solar LLM이\n한 편의 취향 리포트로 해석"),
    ("🧹", "구독 정리 코치", "휴면·저관심 채널의 해제/보류를\nLLM이 사유와 함께 제안"),
    ("✨", "새 채널 추천", "관심사 영상을 역추적해\n실재하는 채널 10개 발굴·추천"),
]
gx, gy = Inches(0.55), Inches(1.7)
cw, ch = Inches(3.95), Inches(2.45)
for i, (ic, h, d) in enumerate(feats):
    col = i % 3; row = i // 3
    l = Emu(int(gx) + col * (int(cw) + int(Inches(0.18))))
    t = Emu(int(gy) + row * (int(ch) + int(Inches(0.22))))
    box(s, l, t, cw, ch, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x30), radius=True)
    text(s, Emu(int(l) + int(Inches(0.25))), Emu(int(t) + int(Inches(0.22))),
         Inches(3.4), Inches(0.6), [[R(ic + "  ", 22, WHITE), R(h, 18, RED, True)]])
    text(s, Emu(int(l) + int(Inches(0.25))), Emu(int(t) + int(Inches(1.0))),
         Inches(3.5), Inches(1.3), [[R(line, 13, GREY)] for line in d.split("\n")],
         line_spacing=1.15)

# ============================================================
# 슬라이드 4 — 시스템 아키텍처
# ============================================================
s = slide()
title(s, "03 · ARCHITECTURE", "시스템 아키텍처")
# 3개 컬럼: 브라우저 / Flask / 외부 API
cols = [
    ("브라우저 (Vanilla JS)", WHITE, [
        "GIS 토큰 클라이언트 (OAuth)",
        "YouTube Data API 직접 호출",
        "Takeout 파일 파싱(JS)",
        "결과 렌더링 · 필터 · 요약",
    ]),
    ("Flask 서버 (app.py)", RED, [
        "/api/analyze — 점수·정렬·분류",
        "LRU 캐시 (입력 의존 키)",
        "Heap Sort · Binary Search",
        "LLM 프록시 7종 엔드포인트",
    ]),
    ("외부 서비스", GREEN, [
        "YouTube Data API v3",
        "  subscriptions · videos",
        "  search · commentThreads",
        "Upstage Solar (solar-pro2)",
    ]),
]
cw = Inches(3.95)
gx = Inches(0.55)
for i, (h, c, items) in enumerate(cols):
    l = Emu(int(gx) + i * (int(cw) + int(Inches(0.18))))
    box(s, l, Inches(1.75), cw, Inches(3.6), fill=PANEL,
        line=RGBColor(0x2A, 0x2A, 0x30), radius=True)
    box(s, l, Inches(1.75), cw, Inches(0.62), fill=RGBColor(0x20, 0x20, 0x26), radius=True)
    text(s, l, Inches(1.75), cw, Inches(0.62), [[R(h, 15, c, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(l) + int(Inches(0.28))), Inches(2.6), Inches(3.45), Inches(2.6),
         [[R("• " + it if not it.startswith("  ") else it, 13,
             GREY if it.startswith("  ") else WHITE)] for it in items],
         line_spacing=1.25)
    if i < 2:
        ax = Emu(int(l) + int(cw) - int(Inches(0.02)))
        text(s, ax, Inches(3.2), Inches(0.22), Inches(0.5),
             [[R("→", 22, RED, True)]], align=PP_ALIGN.CENTER)
# 하단 데이터 흐름
box(s, Inches(0.55), Inches(5.6), Inches(12.23), Inches(1.25), fill=RGBColor(0x12, 0x12, 0x15), radius=True)
text(s, Inches(0.8), Inches(5.72), Inches(11.8), Inches(0.4),
     [[R("데이터 흐름", 13, RED, True, MONO)]])
text(s, Inches(0.8), Inches(6.12), Inches(11.9), Inches(0.7),
     [[R("OAuth 로그인 → 구독·좋아요·댓글 수집 → (선택) Takeout 보강 → ", 13, GREY),
       R("/api/analyze", 13, WHITE, True, MONO),
       R(" 점수 파이프라인 → 그룹 분류·정렬 → 화면 렌더 → ", 13, GREY),
       R("AI 분석 도구(LLM)", 13, GREEN, True)]], line_spacing=1.2)

# ============================================================
# 슬라이드 5 — 사용 흐름 (스크린샷 2장)
# ============================================================
s = slide()
title(s, "04 · USER FLOW", "사용 흐름 — 4단계")
steps = "설정 완료  →  Google 로그인  →  데이터 수집  →  분석 완료"
text(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(0.4),
     [[R(steps, 15, GREEN, True, MONO)]], align=PP_ALIGN.CENTER)
add_image_fit(s, imgpath("기본화면.png"), Inches(0.55), Inches(2.0), Inches(6.0), Inches(4.5))
add_image_fit(s, imgpath("로그인 후.png"), Inches(6.78), Inches(2.0), Inches(6.0), Inches(4.5))
text(s, Inches(0.55), Inches(6.6), Inches(6.0), Inches(0.5),
     [[R("① 클라이언트 ID 설정 후 Google 계정 로그인 (youtube.readonly)", 12, GREY)]],
     align=PP_ALIGN.CENTER)
text(s, Inches(6.78), Inches(6.6), Inches(6.0), Inches(0.5),
     [[R("② (선택) Takeout 업로드 → 댓글·시청 기록 정확 반영", 12, GREY)]],
     align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 6 — 점수 산정 파이프라인
# ============================================================
s = slide()
title(s, "05 · ALGORITHM", "점수 산정 — 5단계 파이프라인")
pipe = [
    ("1", "활동 가중치 합산", "댓글×5 + 시청×3 + 좋아요×2\n+ 장기구독 보너스(최대 +12)", "활동 0 → 즉시 0점"),
    ("2", "Recency Decay", "score × e^(−0.005·d)\nd = 마지막 시청 후 경과일", "오래 안 본 채널 감쇠"),
    ("3", "Min-Max 정규화", "(s−min)/(max−min)×100\n→ 0~100 상대 평가", "사용자 내 비교"),
    ("4", "Binary Search 분류", "경계 [20·50·80] 이진 탐색\n→ 4그룹 매핑 O(log k)", "🔥 😊 👋 🌱"),
    ("5", "Heap Sort 정렬", "직접 구현한 max-heap\n점수 내림차순 O(n log n)", "순위표 출력"),
]
cw = Inches(2.35); gx = Inches(0.55)
for i, (n, h, body, tag) in enumerate(pipe):
    l = Emu(int(gx) + i * (int(cw) + int(Inches(0.07))))
    box(s, l, Inches(1.95), cw, Inches(3.6), fill=PANEL,
        line=RGBColor(0x2A, 0x2A, 0x30), radius=True)
    cir = box(s, Emu(int(l) + int(Inches(0.85))), Inches(2.15), Inches(0.62), Inches(0.62),
              fill=RED, radius=True)
    text(s, Emu(int(l) + int(Inches(0.85))), Inches(2.15), Inches(0.62), Inches(0.62),
         [[R(n, 20, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, l, Inches(2.95), cw, Inches(0.5), [[R(h, 14, WHITE, True)]],
         align=PP_ALIGN.CENTER)
    text(s, Emu(int(l) + int(Inches(0.15))), Inches(3.55), Emu(int(cw) - int(Inches(0.3))),
         Inches(1.3), [[R(ln, 11, GREY, False, MONO)] for ln in body.split("\n")],
         align=PP_ALIGN.CENTER, line_spacing=1.15)
    box(s, Emu(int(l) + int(Inches(0.2))), Inches(4.85),
        Emu(int(cw) - int(Inches(0.4))), Inches(0.45), fill=RGBColor(0x2A, 0x1A, 0x1A),
        radius=True)
    text(s, Emu(int(l) + int(Inches(0.2))), Inches(4.85),
         Emu(int(cw) - int(Inches(0.4))), Inches(0.45), [[R(tag, 11, ORANGE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        text(s, Emu(int(l) + int(cw) - int(Inches(0.05))), Inches(3.4),
             Inches(0.2), Inches(0.4), [[R("→", 16, RED, True)]], align=PP_ALIGN.CENTER)
# 하단: LRU 캐시
box(s, Inches(0.55), Inches(5.75), Inches(12.23), Inches(1.05),
    fill=RGBColor(0x12, 0x12, 0x15), radius=True)
text(s, Inches(0.8), Inches(5.85), Inches(11.8), Inches(0.9),
     [[R("+ LRU 캐시", 14, GREEN, True), R("  (OrderedDict, 입력 의존 키)", 12, GREY, False, MONO)],
      [R("캐시 키 = (id, comments, likes, watchCount, subMonths, lastWatchDays) → 입력이 바뀌면 자동 무효화, 반복 호출은 즉시 히트",
         12, GREY)]], line_spacing=1.2)

# ============================================================
# 슬라이드 7 — 가중치 & 검증 사례
# ============================================================
s = slide()
title(s, "05 · ALGORITHM", "왜 \"활동 0 → 0점\"인가")
# 좌: 가중치
box(s, Inches(0.55), Inches(1.7), Inches(5.6), Inches(2.6), fill=PANEL, radius=True)
text(s, Inches(0.8), Inches(1.85), Inches(5.0), Inches(0.4), [[R("신호별 가중치", 16, RED, True)]])
wt = [("댓글 수", "×5", "글쓰기 = 가장 적극적 교류"),
      ("시청 횟수", "×3", "실제 시청 = 강한 선호"),
      ("좋아요 수", "×2", "클릭 한 번, 중간 신호"),
      ("장기 구독", "+0.5/월", "활동 있을 때만, 2년 캡(+12)")]
for i, (a, b, c) in enumerate(wt):
    y = Emu(int(Inches(2.4)) + i * int(Inches(0.45)))
    text(s, Inches(0.85), y, Inches(1.4), Inches(0.4), [[R(a, 13, WHITE, True)]])
    text(s, Inches(2.25), y, Inches(1.0), Inches(0.4), [[R(b, 13, ORANGE, True, MONO)]])
    text(s, Inches(3.15), y, Inches(2.9), Inches(0.4), [[R(c, 11, GREY)]])
# 좌하: recency 표
box(s, Inches(0.55), Inches(4.45), Inches(5.6), Inches(2.4), fill=PANEL, radius=True)
text(s, Inches(0.8), Inches(4.6), Inches(5.0), Inches(0.4),
     [[R("Recency Decay  ", 16, RED, True), R("e^(−0.005·d)", 13, GREY, False, MONO)]])
dec = [("7일", "0.97"), ("30일", "0.86"), ("90일", "0.64"), ("1년", "0.16"), ("3년", "≈0")]
dx = Inches(0.85)
for d, v in dec:
    box(s, dx, Inches(5.15), Inches(1.0), Inches(1.4), fill=RGBColor(0x1E, 0x1E, 0x24), radius=True)
    text(s, dx, Inches(5.3), Inches(1.0), Inches(0.5), [[R(d, 12, GREY)]], align=PP_ALIGN.CENTER)
    text(s, dx, Inches(5.75), Inches(1.0), Inches(0.6), [[R(v, 16, WHITE, True, MONO)]],
         align=PP_ALIGN.CENTER)
    dx = Emu(int(dx) + int(Inches(1.08)))
# 우: 검증 사례 표
box(s, Inches(6.4), Inches(1.7), Inches(6.38), Inches(5.15), fill=PANEL, radius=True)
text(s, Inches(6.65), Inches(1.85), Inches(5.9), Inches(0.4),
     [[R("검증된 결과 — 동일 입력, 정확한 강등", 15, GREEN, True)]])
rows = [
    ("채널 사례", "댓글", "좋아요", "시청", "점수", WHITE, True),
    ("3일 전 150회 시청", "0", "0", "150", "100", GREEN, False),
    ("5년 활발 (댓글·좋아요)", "20", "80", "0", "60", WHITE, False),
    ("6개월 활발", "5", "20", "0", "15", GREY, False),
    ("10년 잠수 채널", "0", "0", "0", "0", RED, False),
    ("1000일 전 50회 시청", "0", "0", "50", "0", RED, False),
]
ry = Inches(2.4)
colx = [Inches(6.65), Inches(9.55), Inches(10.45), Inches(11.3), Inches(12.1)]
for ri, (name, cm, lk, wc, sc, col, hd) in enumerate(rows):
    y = Emu(int(ry) + ri * int(Inches(0.62)))
    if hd:
        box(s, Inches(6.55), y, Inches(6.05), Inches(0.55), fill=RGBColor(0x24, 0x24, 0x2A))
    elif ri % 2 == 0:
        box(s, Inches(6.55), y, Inches(6.05), Inches(0.55), fill=RGBColor(0x18, 0x18, 0x1D))
    sz = 13 if hd else 13
    text(s, colx[0], y, Inches(2.8), Inches(0.55), [[R(name, sz, col, hd or sc != "")]],
         anchor=MSO_ANCHOR.MIDDLE)
    for ci, val in enumerate([cm, lk, wc]):
        text(s, colx[ci + 1], y, Inches(0.85), Inches(0.55), [[R(val, 13, GREY, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, colx[4], y, Inches(0.6), Inches(0.55),
         [[R(sc, 15 if not hd else 13, col, True, MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(6.65), Inches(6.25), Inches(5.9), Inches(0.55),
     [[R("→ 이전 산식의 함정(", 12, GREY), R("10년 잠수 > 6개월 활발", 12, RED, True),
       R(")을 \"활동 0 → 0점\"으로 해결", 12, GREY)]], line_spacing=1.1)

# ============================================================
# 슬라이드 8 — 적용 알고리즘 목록 (강의 개념)
# ============================================================
s = slide()
title(s, "06 · COURSE CONCEPTS", "코드에 실제로 쓰인 알고리즘")
algos = [
    ("Heap Sort", "재귀 heapify, max-heap", "O(n log n)", "점수 내림차순 정렬"),
    ("Binary Search", "경계 배열 이진 탐색", "O(log k)", "점수 → 4그룹 매핑"),
    ("LRU Cache", "OrderedDict 입력 의존 키", "O(1)", "점수 캐싱·자동 무효화"),
    ("Hash Map", "likeMap·commentMap 등", "O(1) avg", "다채널 집계"),
    ("Min-Max 정규화", "선형 스케일링", "O(n)", "0~100 점수화"),
    ("Recency Decay", "지수 감쇠 exp(−λ·d)", "O(1)", "시청 최근성 반영"),
    ("Cursor 페이지네이션", "nextPageToken 순회", "O(n)", "구독·좋아요 전량 수집"),
    ("병렬 배치 + 백오프", "Promise.all 10개·429 재시도", "O(n/p)", "API 동시 수집"),
]
cw = Inches(6.0); ch = Inches(1.15); gx = Inches(0.55)
for i, (name, how, cx, role) in enumerate(algos):
    col = i % 2; row = i // 2
    l = Emu(int(gx) + col * (int(cw) + int(Inches(0.23))))
    t = Emu(int(Inches(1.7)) + row * (int(ch) + int(Inches(0.12))))
    box(s, l, t, cw, ch, fill=PANEL, line=RGBColor(0x2A, 0x2A, 0x30), radius=True)
    text(s, Emu(int(l) + int(Inches(0.25))), Emu(int(t) + int(Inches(0.12))),
         Inches(3.6), Inches(0.4), [[R(name, 16, RED, True)]])
    text(s, Emu(int(l) + int(Inches(0.25))), Emu(int(t) + int(Inches(0.55))),
         Inches(4.2), Inches(0.5), [[R(how, 12, GREY)]])
    box(s, Emu(int(l) + int(cw) - int(Inches(1.65))), Emu(int(t) + int(Inches(0.16))),
        Inches(1.45), Inches(0.4), fill=RGBColor(0x1E, 0x2A, 0x1E), radius=True)
    text(s, Emu(int(l) + int(cw) - int(Inches(1.65))), Emu(int(t) + int(Inches(0.16))),
         Inches(1.45), Inches(0.4), [[R(cx, 12, GREEN, True, MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(l) + int(cw) - int(Inches(2.9))), Emu(int(t) + int(Inches(0.62))),
         Inches(2.65), Inches(0.4), [[R(role, 11, DIM)]], align=PP_ALIGN.RIGHT)

# ============================================================
# 슬라이드 9 — 분석 결과 화면 (스크린샷)
# ============================================================
s = slide()
title(s, "07 · RESULT", "분석 결과 — 관심도 순위 + AI 도구")
add_image_fit(s, imgpath("우하단.png"), Inches(0.55), Inches(1.6), Inches(8.1), Inches(5.4))
box(s, Inches(8.85), Inches(1.6), Inches(3.93), Inches(5.4), fill=PANEL, radius=True)
text(s, Inches(9.1), Inches(1.85), Inches(3.5), Inches(0.5), [[R("화면 구성", 16, RED, True)]])
text(s, Inches(9.1), Inches(2.45), Inches(3.55), Inches(4.4),
     [[R("• 채널별 ", 13, WHITE), R("관심도 점수·순위", 13, WHITE, True)],
      [R("  (#순위 · 댓글 · 좋아요 · 구독)", 11, GREY)],
      [R("", 7, GREY)],
      [R("• 점수 0 채널 = 휴면으로", 13, WHITE)],
      [R("  하단에 정렬", 13, GREY)],
      [R("", 7, GREY)],
      [R("• 하단 ", 13, WHITE), R("🤖 AI 분석 도구", 13, GREEN, True)],
      [R("  3개 모드 선택", 13, GREY)],
      [R("   - 시청 취향 페르소나", 12, GREY)],
      [R("   - 구독 정리 코치", 12, GREY)],
      [R("   - 새 채널 추천", 12, GREY)],
      [R("", 7, GREY)],
      [R("• 우하단 플로팅 버튼으로", 13, WHITE)],
      [R("  언제든 재호출", 13, GREY)]], line_spacing=1.12)

# ============================================================
# 슬라이드 10 — 페르소나 (스크린샷)
# ============================================================
s = slide()
title(s, "08 · AI TOOL", "시청 취향 페르소나")
add_image_fit(s, imgpath("페르소나.png"), Inches(0.55), Inches(1.6), Inches(8.1), Inches(5.4))
box(s, Inches(8.85), Inches(1.6), Inches(3.93), Inches(5.4), fill=PANEL, radius=True)
text(s, Inches(9.1), Inches(1.85), Inches(3.5), Inches(0.5),
     [[R("Solar LLM 해석", 16, RED, True)]])
text(s, Inches(9.1), Inches(2.45), Inches(3.55), Inches(4.4),
     [[R("채널 점수·카테고리 분포를", 13, WHITE)],
      [R("입력으로 받아 한 편의", 13, GREY)],
      [R("취향 리포트로 풀어냄", 13, GREY)],
      [R("", 8, GREY)],
      [R("예시 결과:", 13, ORANGE, True)],
      [R("\"자취 요리 + 게임/역사 덕후\"", 13, WHITE, True)],
      [R("", 6, GREY)],
      [R("• 자취 요리 마스터", 12, GREY)],
      [R("• 게임 전략가", 12, GREY)],
      [R("• 역사 + 신조어 덕후", 12, GREY)],
      [R("• 건강식 실험가", 12, GREY)],
      [R("", 8, GREY)],
      [R("→ 태그·소제목·요약문을", 12, DIM)],
      [R("   구조화 JSON으로 생성", 12, DIM)]], line_spacing=1.12)

# ============================================================
# 슬라이드 11 — 구독 정리 코치 (스크린샷)
# ============================================================
s = slide()
title(s, "09 · AI TOOL", "구독 정리 코치")
add_image_fit(s, imgpath("구독 정리 코치.png"), Inches(0.55), Inches(1.6), Inches(8.1), Inches(5.4))
box(s, Inches(8.85), Inches(1.6), Inches(3.93), Inches(5.4), fill=PANEL, radius=True)
text(s, Inches(9.1), Inches(1.85), Inches(3.5), Inches(0.5), [[R("정리 제안 로직", 16, RED, True)]])
text(s, Inches(9.1), Inches(2.45), Inches(3.55), Inches(4.4),
     [[R("점수 20 미만 휴면·저관심", 13, WHITE)],
      [R("채널을 LLM이 판정", 13, GREY)],
      [R("", 8, GREY)],
      [R("🔴 ", 13, WHITE), R("해제 추천", 13, RED, True)],
      [R("   장기 미활동(2년+)", 12, GREY)],
      [R("", 6, GREY)],
      [R("🟡 ", 13, WHITE), R("보류 권장", 13, ORANGE, True)],
      [R("   1년 내 미활동", 12, GREY)],
      [R("", 8, GREY)],
      [R("각 채널마다 사유 명시", 13, WHITE)],
      [R("(\"1530일 전 활동 추정,", 12, GREY)],
      [R("  장기 미활동 51개월\")", 12, GREY)]], line_spacing=1.12)

# ============================================================
# 슬라이드 12 — 새 채널 추천 + video-first
# ============================================================
s = slide()
title(s, "10 · AI TOOL", "새 채널 추천 — 영상 우선 발굴")
add_image_fit(s, imgpath("채널 추천.png"), Inches(0.55), Inches(1.6), Inches(7.0), Inches(5.4))
box(s, Inches(7.75), Inches(1.6), Inches(5.03), Inches(5.4), fill=PANEL, radius=True)
text(s, Inches(8.0), Inches(1.8), Inches(4.6), Inches(0.6),
     [[R("환각 방지: 이름을 짓지 않고", 15, RED, True)],
      [R("실재 영상을 역추적", 15, RED, True)]])
disc = [
    ("①", "관심사 키워드 추출", "상위 채널 + 최근 영상 제목"),
    ("②", "키워드별 영상 검색", "search.list → channelId 수집"),
    ("③", "channelId 빈도 집계", "2개+ 키워드 등장 = 진짜 주제"),
    ("④", "후보 프로파일링", "실제 설명·영상 제목 수집 (1 unit)"),
    ("⑤", "LLM 큐레이션", "순위·이유·적합성(fit) 판정"),
]
for i, (n, h, d) in enumerate(disc):
    y = Emu(int(Inches(2.7)) + i * int(Inches(0.82)))
    text(s, Inches(8.0), y, Inches(0.5), Inches(0.6), [[R(n, 18, GREEN, True)]])
    text(s, Inches(8.5), y, Inches(4.1), Inches(0.4), [[R(h, 13, WHITE, True)]])
    text(s, Inches(8.5), Emu(int(y) + int(Inches(0.33))), Inches(4.1), Inches(0.4),
         [[R(d, 11, GREY, False, MONO)]])
text(s, Inches(8.0), Inches(6.85), Inches(4.6), Inches(0.4),
     [[R("비싼 search(100u)는 ②에서만, 깊은 분석은 1u로 → 쿼터 절약", 10, DIM)]])

# ============================================================
# 슬라이드 13 — Takeout 하이브리드 전략
# ============================================================
s = slide()
title(s, "11 · DATA STRATEGY", "Takeout 하이브리드 — 비용 vs 정확도")
text(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(0.5),
     [[R("YouTube API엔 \"내 댓글·시청 기록 전체\" 조회가 없다 → ", 14, GREY),
       R("Google Takeout 파일을 보조 입력으로 받아 보강", 14, WHITE, True)]])
hdr = ["모드", "사용자 부담", "API 비용", "정확도"]
data = [
    ("둘 다 미업로드", "없음", "commentThreads × 채널 수", "댓글 근사 · 시청 미반영", GREY),
    ("comments.csv 만", "1회 업로드", "videos.list × 영상수/50", "댓글 정확 · 시청 미반영", WHITE),
    ("watch-history.json 만", "1회 업로드", "0 units (subtitles에 ID)", "시청 정확 · 댓글 근사", WHITE),
    ("둘 다 ⭐", "2회 업로드", "댓글 영상 해석만", "모두 정확", GREEN),
]
colx = [Inches(0.55), Inches(3.7), Inches(6.0), Inches(9.3)]
colw = [Inches(3.15), Inches(2.3), Inches(3.3), Inches(3.45)]
# 헤더
box(s, Inches(0.55), Inches(2.2), Inches(12.23), Inches(0.55), fill=RGBColor(0x24, 0x24, 0x2A))
for ci, h in enumerate(hdr):
    text(s, colx[ci], Inches(2.2), colw[ci], Inches(0.55), [[R(h, 13, RED, True)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
for ri, (m, u, c, a, col) in enumerate(data):
    y = Emu(int(Inches(2.78)) + ri * int(Inches(0.78)))
    if ri % 2 == 0:
        box(s, Inches(0.55), y, Inches(12.23), Inches(0.72), fill=RGBColor(0x16, 0x16, 0x1A))
    if col == GREEN:
        box(s, Inches(0.55), y, Inches(12.23), Inches(0.72), fill=RGBColor(0x14, 0x22, 0x18))
    text(s, colx[0], y, colw[0], Inches(0.72), [[R(m, 14, col, col in (GREEN, WHITE))]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, colx[1], y, colw[1], Inches(0.72), [[R(u, 12, GREY)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text(s, colx[2], y, colw[2], Inches(0.72), [[R(c, 12, GREY, False, MONO)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text(s, colx[3], y, colw[3], Inches(0.72), [[R(a, 12, col)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
text(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.9),
     [[R("자동 파일 식별: ", 13, ORANGE, True),
       R("파일명(watch/comment/시청/댓글) + 내용 첫 4KB(titleUrl·subtitles 패턴)로 종류 자동 분기", 13, GREY)],
      [R("watch-history.json은 subtitles[0].url에 channelId가 박혀 있어 ", 12, GREY),
       R("API 호출 0회", 12, GREEN, True), R("로 시청 횟수·마지막 시청일 추출", 12, GREY)]],
     line_spacing=1.2)

# ============================================================
# 슬라이드 14 — 제약 & 향후
# ============================================================
s = slide()
title(s, "12 · LIMITS & NEXT", "제약 사항과 향후 과제")
box(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.1), fill=PANEL, radius=True)
text(s, Inches(0.85), Inches(1.9), Inches(5.4), Inches(0.4), [[R("현재 제약", 17, RED, True)]])
text(s, Inches(0.85), Inches(2.5), Inches(5.45), Inches(4.2),
     [[R("• YouTube API 일일 쿼터 ", 14, WHITE), R("10,000 units", 14, ORANGE, True)],
      [R("  단일 GCP 프로젝트를 모든 사용자 공유", 12, GREY)],
      [R("", 7, GREY)],
      [R("• OAuth access_token ", 14, WHITE), R("1시간 만료", 14, ORANGE, True)],
      [R("  갱신 로직 없음 (단발성 토큰)", 12, GREY)],
      [R("", 7, GREY)],
      [R("• 댓글 API 근사: 채널당 최근 100개만", 14, WHITE)],
      [R("  스캔 → Takeout으로 보완", 12, GREY)],
      [R("", 7, GREY)],
      [R("• 점수는 ", 14, WHITE), R("상대 평가", 14, ORANGE, True),
       R(" — 사용자 간 비교 불가", 14, GREY)]], line_spacing=1.12)
box(s, Inches(6.78), Inches(1.7), Inches(6.0), Inches(5.1), fill=PANEL, radius=True)
text(s, Inches(7.08), Inches(1.9), Inches(5.4), Inches(0.4), [[R("향후 과제", 17, GREEN, True)]])
text(s, Inches(7.08), Inches(2.5), Inches(5.45), Inches(4.2),
     [[R("✓ Refresh Token 도입 → 장시간 분석", 14, WHITE)],
      [R("", 7, GREY)],
      [R("✓ 분석 결과 저장·기간별 추이 비교", 14, WHITE)],
      [R("", 7, GREY)],
      [R("✓ 추천 채널 원클릭 구독 연동", 14, WHITE)],
      [R("", 7, GREY)],
      [R("✓ 카테고리 군집화 시각화(차트)", 14, WHITE)],
      [R("", 7, GREY)],
      [R("✓ OAuth 동의 화면 프로덕션 게시", 14, WHITE)],
      [R("  → 테스트 사용자 100명 제한 해제", 12, GREY)]], line_spacing=1.12)

# ============================================================
# 슬라이드 15 — 마무리
# ============================================================
s = slide()
box(s, Inches(0), Inches(0), SW, SH, fill=BG)
box(s, Inches(0), Inches(0), Inches(0.13), SH, fill=RED)
text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
     [[R("// THANK YOU", 18, RED, True, MONO)]])
text(s, Inches(0.85), Inches(2.85), Inches(11.7), Inches(1.0),
     [[R("구독 채널 ", 46, WHITE, True), R("관심도 분석기", 46, RED, True)]])
text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.5),
     [[R("점수화 알고리즘 + Solar LLM 해석으로", 18, GREY)],
      [R("\"내가 진짜 보는 채널\"을 찾아주는 웹 애플리케이션", 18, GREY)]], line_spacing=1.25)
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.5),
     [[R("Flask · Vanilla JS · YouTube Data API v3 · Upstage Solar", 14, DIM, False, MONO)]])
text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
     [[R("Q & A", 22, WHITE, True)]])

out = os.path.join(os.path.dirname(__file__), "구독채널_관심도_분석기_발표.pptx")
prs.save(out)
print("저장 완료:", out, "| 슬라이드", len(prs.slides._sldIdLst))
